#!/usr/bin/env python3
"""Generate presentation.pptx from the Remotion + Claude Code HTML slideshow."""

import math
import os
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Constants ──────────────────────────────────────────────────────────────

BG_COLOR = RGBColor(0x1B, 0x1B, 0x1B)
HEADING_COLOR = RGBColor(0xDC, 0xDA, 0xD5)
BODY_COLOR = RGBColor(0xA0, 0x9D, 0x96)
MUTED_COLOR = RGBColor(0xC5, 0xC1, 0xB9)
ACCENT = RGBColor(0x22, 0xC5, 0x5E)
PANEL_BG = RGBColor(0x25, 0x25, 0x25)
DARK_PANEL = RGBColor(0x1E, 0x1E, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DOT_RED = RGBColor(0xFF, 0x5F, 0x57)
DOT_YELLOW = RGBColor(0xFE, 0xBC, 0x2E)
DOT_GREEN_DOT = RGBColor(0x28, 0xC8, 0x40)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_HEBREW = "Heebo"
FONT_MONO = "Consolas"


# ── Helpers ────────────────────────────────────────────────────────────────

BG_IMG_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_slide_bg.png")
GRID_IMG_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_slide_grid.png")


def _generate_bg_images():
    """Generate PNGs: base bg with green blob, and separate grid overlay."""
    if not os.path.exists(BG_IMG_PATH):
        w, h = 1920, 1080
        img = Image.new("RGB", (w, h), (0x1B, 0x1B, 0x1B))

        # ── Green radial blob ──
        blob = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        cx, cy = w // 2, h // 2
        max_r = int(math.hypot(cx, cy) * 0.7)
        for r in range(max_r, 0, -1):
            t = r / max_r
            if t < 0.0:
                alpha = int(255 * 0.20)
            elif t < 0.4:
                frac = t / 0.4
                alpha = int(255 * (0.20 + (0.06 - 0.20) * frac))
            elif t < 1.0:
                frac = (t - 0.4) / 0.6
                alpha = int(255 * (0.06 * (1 - frac)))
            else:
                alpha = 0
            if alpha <= 0:
                continue
            draw = ImageDraw.Draw(blob)
            draw.ellipse([cx - r, cy - r, cx + r, cy + r],
                         fill=(0x22, 0xC5, 0x5E, alpha))

        img = Image.alpha_composite(img.convert("RGBA"), blob).convert("RGB")
        img.save(BG_IMG_PATH, "PNG")

    if not os.path.exists(GRID_IMG_PATH):
        # Grid image is 20% larger than slide so the drift animation
        # doesn't reveal blank edges
        gw, gh = 2304, 1296
        grid = Image.new("RGBA", (gw, gh), (0, 0, 0, 0))
        draw = ImageDraw.Draw(grid)
        cell_size = 40
        line_alpha = int(255 * 0.09)
        grid_col = (255, 255, 255, line_alpha)
        for x in range(0, gw, cell_size):
            draw.line([(x, 0), (x, gh)], fill=grid_col, width=1)
        for y in range(0, gh, cell_size):
            draw.line([(0, y), (gw, y)], fill=grid_col, width=1)
        grid.save(GRID_IMG_PATH, "PNG")


def _add_grid_drift_animation(slide, pic_element):
    """Add a slow infinite diagonal drift motion-path animation to the grid."""
    # We need to add animation XML to the slide's timing tree.
    # The motion path moves the grid gently so it appears alive.
    sld = slide._element
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ans = "http://schemas.openxmlformats.org/drawingml/2006/main"
    pns = "http://schemas.microsoft.com/office/powerpoint/2010/main"

    nv = pic_element.find(qn("p:nvPicPr"))
    sp_id = nv.find(qn("p:cNvPr")).get("id")

    # Build timing XML for an infinite, auto-starting motion path
    timing_xml = f'''<p:timing xmlns:p="{ns}"
        xmlns:a="{ans}">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                          <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                          <p:par>
                            <p:cTn id="4" fill="hold">
                              <p:stCondLst>
                                <p:cond delay="0"/>
                              </p:stCondLst>
                              <p:childTnLst>
                                <p:animMotion origin="layout" path="M 0 0 L 0.03 0.03" pathEditMode="relative"
                                    ptsTypes="AA" rAng="0">
                                  <p:cBhvr>
                                    <p:cTn id="5" dur="8000" repeatCount="indefinite" autoRev="1" fill="hold">
                                      <p:stCondLst>
                                        <p:cond delay="0"/>
                                      </p:stCondLst>
                                    </p:cTn>
                                    <p:tgtEl>
                                      <p:spTgt spid="{sp_id}"/>
                                    </p:tgtEl>
                                  </p:cBhvr>
                                </p:animMotion>
                              </p:childTnLst>
                            </p:cTn>
                          </p:par>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst>
                  <p:cond evt="onPrev" delay="0">
                    <p:tgtEl><p:sldTgt/></p:tgtEl>
                  </p:cond>
                </p:prevCondLst>
                <p:nextCondLst>
                  <p:cond evt="onNext" delay="0">
                    <p:tgtEl><p:sldTgt/></p:tgtEl>
                  </p:cond>
                </p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
    </p:timing>'''

    from lxml import etree
    timing_el = etree.fromstring(timing_xml)
    # Remove any existing timing
    existing = sld.find(qn("p:timing"))
    if existing is not None:
        sld.remove(existing)
    sld.append(timing_el)


def set_slide_bg(slide, color=BG_COLOR):
    """Set slide background to the pre-generated PNG images with animated grid."""
    _generate_bg_images()
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

    sp_tree = slide.shapes._spTree

    # Base bg image (dark + green blob)
    pic = slide.shapes.add_picture(BG_IMG_PATH, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)

    # Grid overlay — slightly oversized, centred, so drift doesn't show edges
    overflow = Inches(1.5)
    grid_pic = slide.shapes.add_picture(
        GRID_IMG_PATH,
        Emu(int(-overflow / 2)), Emu(int(-overflow / 2)),
        Emu(int(SLIDE_W + overflow)), Emu(int(SLIDE_H + overflow))
    )
    sp_tree.remove(grid_pic._element)
    sp_tree.insert(3, grid_pic._element)

    # Animate the grid with a slow diagonal drift
    _add_grid_drift_animation(slide, grid_pic._element)


def _set_rtl(paragraph, rtl=True):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1" if rtl else "0")


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=HEADING_COLOR, bold=False, alignment=PP_ALIGN.RIGHT,
                rtl=True, font_name=FONT_HEBREW, line_spacing=None,
                word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    _set_rtl(p, rtl)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox, tf, p


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=18,
                          color=HEADING_COLOR, bold=False, alignment=PP_ALIGN.RIGHT,
                          rtl=True, font_name=FONT_HEBREW, line_spacing=None):
    """Add textbox with multiple lines, each as a separate paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        _set_rtl(p, rtl)
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox, tf


def add_slide_label(slide, text, top=Inches(0.6)):
    """Add green uppercase label at top center."""
    add_textbox(slide, Inches(0), top, SLIDE_W, Inches(0.5), text,
                font_size=14, color=ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER, rtl=False, font_name=FONT_HEBREW)


def add_rounded_rect(slide, left, top, width, height, fill_color=PANEL_BG,
                     border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    # Reduce corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_terminal_dots(slide, left, top):
    """Add red/yellow/green dots for terminal bar."""
    colors = [DOT_RED, DOT_YELLOW, DOT_GREEN_DOT]
    for i, c in enumerate(colors):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     left + Inches(i * 0.28), top,
                                     Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()


def add_pill(slide, left, top, width, height, text, fill_color=None,
             border_color=ACCENT, text_color=WHITE, font_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.5  # fully rounded ends
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    run.font.name = FONT_HEBREW
    run.font.bold = True
    return shape


def add_code_block(slide, left, top, width, height, lines,
                   font_size=11, title=None):
    """Add a terminal-style code block with dots and monospace text."""
    bar_h = Inches(0.35)
    add_rounded_rect(slide, left, top, width, height, fill_color=DARK_PANEL,
                     border_color=RGBColor(0x33, 0x33, 0x33))
    # dots
    add_terminal_dots(slide, left + Inches(0.15), top + Inches(0.1))
    # title
    if title:
        add_textbox(slide, left + Inches(1.2), top + Inches(0.05),
                    width - Inches(1.5), Inches(0.3), title,
                    font_size=9, color=RGBColor(0x66, 0x66, 0x66),
                    alignment=PP_ALIGN.LEFT, rtl=False, font_name=FONT_MONO)
    # code lines
    code_text = "\n".join(lines)
    add_textbox(slide, left + Inches(0.2), top + bar_h,
                width - Inches(0.4), height - bar_h - Inches(0.1),
                code_text, font_size=font_size, color=MUTED_COLOR,
                alignment=PP_ALIGN.LEFT, rtl=False, font_name=FONT_MONO)


# ── Slide Builders ─────────────────────────────────────────────────────────

def build_slide_1_title(prs):
    """Slide 1: Opening - Remotion + Claude Code"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Green tag
    add_textbox(slide, Inches(0), Inches(2.0), SLIDE_W, Inches(0.5),
                "FROM CODE TO VIDEO", font_size=14, color=ACCENT,
                bold=True, alignment=PP_ALIGN.CENTER, rtl=False)

    # Main title
    add_textbox(slide, Inches(1.5), Inches(2.7), Inches(10.3), Inches(1.2),
                "\u05de\u05d4\u05e4\u05db\u05ea \u05d9\u05e6\u05d9\u05e8\u05ea \u05d4\u05d5\u05d5\u05d9\u05d3\u05d0\u05d5 \u05e2\u05dd Remotion \u05d5-Claude Code",
                font_size=36, color=HEADING_COLOR, bold=True,
                alignment=PP_ALIGN.CENTER, rtl=True)

    # Subtitle
    add_multiline_textbox(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(1.0),
                          ["\u05d0\u05e0\u05d9\u05de\u05e6\u05d9\u05d4. \u05d2\u05e8\u05e4\u05d9\u05e7\u05d4. \u05e1\u05d0\u05d5\u05e0\u05d3. \u05e7\u05e8\u05d9\u05d9\u05e0\u05d5\u05ea.", "\u05d4\u05db\u05dc \u05e0\u05d1\u05e0\u05d4 \u05d1\u05e7\u05d5\u05d3."],
                          font_size=20, color=BODY_COLOR,
                          alignment=PP_ALIGN.CENTER, rtl=True, line_spacing=32)


def build_slide_2_code_to_video(prs):
    """Slide 2: Code -> Video visual (LTR)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Code panel (left)
    code_lines = [
        "import { Composition } from 'remotion'",
        "import { useCurrentFrame } from 'remotion'",
        "import { interpolate } from 'remotion'",
        "",
        "export const MyVideo = () => {",
        "  const frame = useCurrentFrame()",
        "  const opacity = interpolate(",
        "    frame, [0, 30], [0, 1]",
        "  )",
        "  return <AbsoluteFill style={{opacity}} />",
        "}",
    ]
    add_code_block(slide, Inches(0.8), Inches(1.0), Inches(5.2), Inches(5.5),
                   code_lines, font_size=12, title="MyVideo.tsx")

    # Arrow connector
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(6.3), Inches(3.2),
                                   Inches(0.8), Inches(0.6))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT
    arrow.line.fill.background()

    # Video panel (right)
    panel_l = Inches(7.4)
    panel_t = Inches(1.0)
    panel_w = Inches(5.2)
    panel_h = Inches(5.5)
    add_rounded_rect(slide, panel_l, panel_t, panel_w, panel_h,
                     fill_color=DARK_PANEL,
                     border_color=RGBColor(0x33, 0x33, 0x33))
    add_terminal_dots(slide, panel_l + Inches(0.15), panel_t + Inches(0.1))

    # Video frame thumbnails
    colors = [
        (0x63, 0x66, 0xF1), (0x22, 0xC5, 0x5E), (0xF5, 0x9E, 0x0B),
        (0x0E, 0xA5, 0xE9), (0xEC, 0x48, 0x99), (0xF9, 0x73, 0x16),
    ]
    frame_w = Inches(1.4)
    frame_h = Inches(1.3)
    start_x = panel_l + Inches(0.4)
    start_y = panel_t + Inches(0.6)
    for i, (r, g, b) in enumerate(colors):
        row = i // 3
        col = i % 3
        x = start_x + col * Inches(1.5)
        y = start_y + row * Inches(1.5)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, y, frame_w, frame_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(r, g, b)
        rect.line.fill.background()
        rect.adjustments[0] = 0.08

    # Timeline bar at bottom of video panel
    add_rounded_rect(slide, panel_l + Inches(0.3),
                     panel_t + Inches(4.8), Inches(4.6), Inches(0.15),
                     fill_color=RGBColor(0x33, 0x33, 0x33))
    # Playhead
    add_rounded_rect(slide, panel_l + Inches(0.3 + 1.5),
                     panel_t + Inches(4.75), Inches(0.08), Inches(0.25),
                     fill_color=ACCENT)


def build_slide_3_problem(prs):
    """Slide 3: The Problem - timeline visual"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Text (right side for RTL)
    text_l = Inches(7.5)
    text_w = Inches(5.3)

    add_textbox(slide, text_l, Inches(1.5), text_w, Inches(0.4),
                "THE PROBLEM", font_size=14, color=ACCENT, bold=True,
                alignment=PP_ALIGN.RIGHT, rtl=False)

    add_textbox(slide, text_l, Inches(2.1), text_w, Inches(0.8),
                "\u05dc\u05e2\u05e8\u05d5\u05da \u05d5\u05d9\u05d3\u05d0\u05d5 \u05d6\u05d4... \u05e7\u05e9\u05d4",
                font_size=34, color=HEADING_COLOR,
                bold=True, alignment=PP_ALIGN.RIGHT, rtl=True)

    add_textbox(slide, text_l, Inches(3.1), text_w, Inches(0.5),
                "\u05d8\u05d9\u05d9\u05de\u05dc\u05d9\u05d9\u05df. \u05dc\u05d9\u05d9\u05e8\u05d9\u05dd. \u05d0\u05e4\u05e7\u05d8\u05d9\u05dd. \u05d0\u05d9\u05e0\u05e1\u05d5\u05e3 \u05e7\u05dc\u05d9\u05e7\u05d9\u05dd.",
                font_size=18, color=ACCENT, alignment=PP_ALIGN.RIGHT, rtl=True)

    add_textbox(slide, text_l, Inches(3.8), text_w, Inches(0.8),
                "\u05e2\u05e8\u05d9\u05db\u05ea \u05d5\u05d9\u05d3\u05d0\u05d5 \u05e7\u05dc\u05d0\u05e1\u05d9\u05ea \u05d4\u05d9\u05d0 \u05d9\u05d3\u05e0\u05d9\u05ea. \u05d7\u05d5\u05d6\u05e8\u05ea \u05e2\u05dc \u05e2\u05e6\u05de\u05d4. \u05d5\u05dc\u05d0 \u05e1\u05e7\u05d9\u05d9\u05dc\u05d1\u05d9\u05dc\u05d9\u05ea.",
                font_size=16, color=BODY_COLOR, alignment=PP_ALIGN.RIGHT, rtl=True)

    # Timeline visual (left side)
    tl_l = Inches(0.5)
    tl_t = Inches(1.2)

    # Ruler ticks
    ticks = ["00:00", "00:05", "00:10", "00:15", "00:20", "00:25", "00:30"]
    for i, tick in enumerate(ticks):
        x = tl_l + Inches(i * 0.9)
        add_textbox(slide, x, tl_t, Inches(0.7), Inches(0.3), tick,
                    font_size=8, color=RGBColor(0x55, 0x55, 0x55),
                    alignment=PP_ALIGN.CENTER, rtl=False, font_name=FONT_MONO)

    # Tracks
    tracks = [
        ("VIDEO",  (0x3B, 0x82, 0xF6), [(0, 22), (25, 27), (54, 18), (76, 20)]),
        ("AUDIO",  (0x22, 0xC5, 0x5E), [(0, 100)]),
        ("TEXT",   (0xF5, 0x9E, 0x0B), [(10, 18), (36, 14), (59, 22)]),
        ("FX",     (0xA8, 0x55, 0xF7), [(7, 11), (32, 9), (52, 13), (71, 10)]),
        ("MUSIC",  (0xEC, 0x48, 0x99), [(0, 42), (47, 30)]),
        ("B-ROLL", (0x06, 0xB6, 0xD4), [(5, 16), (32, 20), (59, 13)]),
    ]

    track_h = Inches(0.45)
    clip_area_l = tl_l + Inches(0.9)
    clip_area_w = Inches(5.5)

    for t_idx, (label, (cr, cg, cb), clips) in enumerate(tracks):
        y = tl_t + Inches(0.4) + t_idx * Inches(0.55)
        # Label
        add_textbox(slide, tl_l, y + Inches(0.05), Inches(0.85), Inches(0.35),
                    label, font_size=8, color=RGBColor(0x77, 0x77, 0x77),
                    alignment=PP_ALIGN.LEFT, rtl=False, font_name=FONT_MONO,
                    bold=True)
        # Clips
        for (start_pct, width_pct) in clips:
            cx = clip_area_l + Emu(int(clip_area_w * start_pct / 100))
            cw = Emu(int(clip_area_w * width_pct / 100))
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          cx, y, cw, track_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(cr, cg, cb)
            rect.line.fill.background()
            rect.adjustments[0] = 0.15
            # Make slightly transparent via XML
            spPr = rect._element.spPr
            solidFill_el = spPr.find(qn("a:solidFill"))
            if solidFill_el is not None:
                srgb = solidFill_el.find(qn("a:srgbClr"))
                if srgb is not None:
                    alpha = srgb.makeelement(qn("a:alpha"), {"val": "60000"})
                    srgb.append(alpha)

    # Playhead line
    ph_x = clip_area_l + Emu(int(clip_area_w * 35 / 100))
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                ph_x, tl_t + Inches(0.3),
                                Inches(0.02), Inches(3.5))
    ph.fill.solid()
    ph.fill.fore_color.rgb = RGBColor(0xEF, 0x44, 0x44)
    ph.line.fill.background()


def build_slide_4_solution(prs):
    """Slide 4: What is Remotion - code + preview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Header
    add_slide_label(slide, "THE SOLUTION")
    add_textbox(slide, Inches(0), Inches(1.0), SLIDE_W, Inches(0.7),
                "\u05de\u05d4 \u05d6\u05d4 Remotion?",
                font_size=32, color=HEADING_COLOR,
                bold=True, alignment=PP_ALIGN.CENTER, rtl=True)
    add_textbox(slide, Inches(0), Inches(1.7), SLIDE_W, Inches(0.4),
                "\u05d1\u05d5\u05e0\u05d9\u05dd \u05d5\u05d9\u05d3\u05d0\u05d5 \u05e2\u05dd React",
                font_size=18, color=ACCENT,
                alignment=PP_ALIGN.CENTER, rtl=True)
    add_textbox(slide, Inches(2), Inches(2.2), Inches(9.3), Inches(0.6),
                "Remotion \u05de\u05d0\u05e4\u05e9\u05e8 \u05dc\u05d1\u05e0\u05d5\u05ea \u05d5\u05d9\u05d3\u05d0\u05d5 \u05d1\u05d0\u05de\u05e6\u05e2\u05d5\u05ea \u05e7\u05d5\u05d3. \u05e7\u05d5\u05de\u05e4\u05d5\u05e0\u05e0\u05d8\u05d5\u05ea. \u05d0\u05e0\u05d9\u05de\u05e6\u05d9\u05d5\u05ea. \u05d3\u05d9\u05e0\u05de\u05d9\u05d5\u05ea \u05de\u05dc\u05d0\u05d4.",
                font_size=14, color=BODY_COLOR, alignment=PP_ALIGN.CENTER, rtl=True)

    # Code panel (left)
    code_lines = [
        "import { AbsoluteFill, useCurrentFrame,",
        "         interpolate } from 'remotion'",
        "",
        "export const MyComp = () => {",
        "  const frame = useCurrentFrame()",
        "  return (<AbsoluteFill style={{",
        "    background: '#0d1117' }}>",
        "    <div style={{ transform:",
        "      `translateX(${frame*3}px)` }} />",
        "    <h1 style={{ opacity:",
        "      interpolate(frame,[0,30],[0,1])",
        "    }}>Hello, Remotion</h1>",
        "    <div style={{ transform:",
        "      `scale(...)` }} />",
        "  </AbsoluteFill>)",
        "}",
    ]
    add_code_block(slide, Inches(0.5), Inches(3.0), Inches(7.0), Inches(4.2),
                   code_lines, font_size=10, title="MyComp.tsx")

    # Preview panel (right)
    prev_l = Inches(8.0)
    prev_t = Inches(3.0)
    prev_w = Inches(4.8)
    prev_h = Inches(4.2)
    add_rounded_rect(slide, prev_l, prev_t, prev_w, prev_h,
                     fill_color=RGBColor(0x0D, 0x11, 0x17),
                     border_color=RGBColor(0x33, 0x33, 0x33))
    add_terminal_dots(slide, prev_l + Inches(0.15), prev_t + Inches(0.1))
    add_textbox(slide, prev_l + Inches(1.2), prev_t + Inches(0.05),
                Inches(2), Inches(0.3), "Preview \u00b7 frame 24",
                font_size=9, color=RGBColor(0x66, 0x66, 0x66),
                alignment=PP_ALIGN.LEFT, rtl=False, font_name=FONT_MONO)

    # Green box
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 prev_l + Inches(0.8), prev_t + Inches(0.8),
                                 Inches(1.2), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT
    box.line.fill.background()

    # "Hello, Remotion" text
    add_textbox(slide, prev_l + Inches(0.5), prev_t + Inches(1.8),
                Inches(3.8), Inches(0.6), "Hello, Remotion",
                font_size=22, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, rtl=False)

    # Green circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    prev_l + Inches(1.8), prev_t + Inches(2.7),
                                    Inches(1.2), Inches(1.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()


def build_slide_5_showcase(prs):
    """Slide 5: 5 use-case cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_slide_label(slide, "\u05de\u05d4 \u05d0\u05e4\u05e9\u05e8 \u05dc\u05d1\u05e0\u05d5\u05ea")
    add_textbox(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(0.8),
                "\u05dc\u05d0 \u05e8\u05e7 \u05d0\u05e0\u05d9\u05de\u05e6\u05d9\u05d4. Remotion \u05d4\u05d9\u05d0 \u05de\u05e2\u05e8\u05db\u05ea \u05dc\u05d9\u05d9\u05e6\u05d5\u05e8 \u05d5\u05d9\u05d3\u05d0\u05d5.",
                font_size=28, color=HEADING_COLOR, bold=True,
                alignment=PP_ALIGN.CENTER, rtl=True)

    # 5 cards
    cards = [
        ("\u05e1\u05e8\u05d8\u05d5\u05e0\u05d9\nSHORTS", "\u25af"),
        ("\u05d4\u05d3\u05de\u05d9\u05d5\u05ea\n\u05e0\u05ea\u05d5\u05e0\u05d9\u05dd", "\u2581\u2583\u2585\u2587"),
        ("\u05e1\u05e8\u05d8\u05d5\u05e0\u05d9\nPromotion", "\u25a0"),
        ("\u05d4\u05d5\u05e1\u05e4\u05ea\n\u05db\u05ea\u05d5\u05d1\u05d9\u05d5\u05ea", "CC"),
        ("\u05d0\u05d9\u05e0\u05d8\u05e8\u05d0\u05e7\u05e6\u05d9\u05d5\u05ea\n\u05d1\u05de\u05de\u05e9\u05e7", "\u25cb\u25cb\u25cb"),
    ]
    card_w = Inches(2.1)
    card_h = Inches(3.5)
    total_w = 5 * card_w + 4 * Inches(0.25)
    start_x = Emu(int((SLIDE_W - total_w) / 2))

    for i, (label, icon) in enumerate(cards):
        x = start_x + i * (card_w + Inches(0.25))
        y = Inches(2.3)

        # Card background
        add_rounded_rect(slide, x, y, card_w, card_h,
                         fill_color=PANEL_BG,
                         border_color=RGBColor(0x33, 0x44, 0x33))

        # Icon area
        add_textbox(slide, x, y + Inches(0.4), card_w, Inches(1.2), icon,
                    font_size=32, color=ACCENT, alignment=PP_ALIGN.CENTER,
                    rtl=False)

        # Label
        lines = label.split("\n")
        add_multiline_textbox(slide, x + Inches(0.1), y + Inches(1.8),
                              card_w - Inches(0.2), Inches(1.2), lines,
                              font_size=14, color=MUTED_COLOR,
                              alignment=PP_ALIGN.CENTER, rtl=True,
                              line_spacing=22)

    # Bottom text
    add_textbox(slide, Inches(0), Inches(6.2), SLIDE_W, Inches(0.5),
                "\u05d5\u05d9\u05d3\u05d0\u05d5 \u05e9\u05de\u05d1\u05d5\u05e1\u05e1 \u05e2\u05dc \u05d3\u05d0\u05d8\u05d4, \u05e4\u05e8\u05d5\u05de\u05e4\u05d8 \u05d0\u05d5 API.",
                font_size=16, color=BODY_COLOR, alignment=PP_ALIGN.CENTER,
                rtl=True)


def build_slide_6_pipeline(prs):
    """Slide 6: Prompt -> Agent -> Code -> Video flow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_slide_label(slide, "\u05dc\u05de\u05d4 \u05e2\u05db\u05e9\u05d9\u05d5")
    add_textbox(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.7),
                "\u05db\u05d9 \u05d1\u05d9\u05e0\u05d4 \u05de\u05dc\u05d0\u05db\u05d5\u05ea\u05d9\u05ea \u05db\u05d5\u05ea\u05d1\u05ea \u05e7\u05d5\u05d3",
                font_size=30, color=HEADING_COLOR,
                bold=True, alignment=PP_ALIGN.CENTER, rtl=True)

    steps = [
        ("PROMPT", "\u05ea\u05d0\u05e8 \u05de\u05d4 \u05d0\u05ea\u05d4 \u05e8\u05d5\u05e6\u05d4"),
        ("AI AGENT", "Claude Code \u05de\u05e0\u05ea\u05d7 \u05d5\u05de\u05ea\u05db\u05e0\u05df"),
        ("CODE", "\u05e7\u05d5\u05de\u05e4\u05d5\u05e0\u05e0\u05d8\u05d5\u05ea React + Remotion"),
        ("VIDEO", "MP4 \u05de\u05d5\u05db\u05df \u05dc\u05e4\u05e8\u05e1\u05d5\u05dd"),
    ]

    box_w = Inches(2.5)
    box_h = Inches(2.8)
    arrow_w = Inches(0.6)
    total = 4 * box_w + 3 * arrow_w
    start_x = Emu(int((SLIDE_W - total) / 2))
    box_y = Inches(2.4)

    for i, (name, desc) in enumerate(steps):
        x = start_x + i * (box_w + arrow_w)

        # Box
        add_rounded_rect(slide, x, box_y, box_w, box_h,
                         fill_color=PANEL_BG,
                         border_color=ACCENT)

        # Step name
        add_textbox(slide, x, box_y + Inches(0.6), box_w, Inches(0.5),
                    name, font_size=20, color=ACCENT, bold=True,
                    alignment=PP_ALIGN.CENTER, rtl=False)

        # Description
        add_textbox(slide, x + Inches(0.2), box_y + Inches(1.3),
                    box_w - Inches(0.4), Inches(1.0),
                    desc, font_size=14, color=BODY_COLOR,
                    alignment=PP_ALIGN.CENTER, rtl=True)

        # Arrow between boxes
        if i < 3:
            ax = x + box_w + Inches(0.05)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           ax, box_y + Inches(1.2),
                                           Inches(0.5), Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    # Bottom note
    add_textbox(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.6),
                "\u05db\u05dc \u05d4\u05e6\u05d9\u05e0\u05d5\u05e8 \u05d4\u05d6\u05d4 \u05e8\u05e5 \u05d0\u05d5\u05d8\u05d5\u05de\u05d8\u05d9\u05ea \u2014 \u05d0\u05ea\u05d4 \u05e8\u05e7 \u05de\u05ea\u05d0\u05e8 \u05d0\u05ea \u05d4\u05d5\u05d9\u05d3\u05d0\u05d5 \u05e9\u05d0\u05ea\u05d4 \u05e8\u05d5\u05e6\u05d4",
                font_size=15, color=BODY_COLOR, alignment=PP_ALIGN.CENTER,
                rtl=True)


def build_slide_7_terminal(prs):
    """Slide 7: Claude Code terminal demo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_slide_label(slide, "\u05de\u05d4 \u05d6\u05d4 Claude Code")
    add_textbox(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.7),
                "\u05e1\u05d5\u05db\u05df AI \u05e9\u05db\u05d5\u05ea\u05d1, \u05de\u05e8\u05d9\u05e5 \u05d5\u05de\u05e8\u05e0\u05d3\u05e8 \u05e7\u05d5\u05d3",
                font_size=30,
                color=HEADING_COLOR, bold=True,
                alignment=PP_ALIGN.CENTER, rtl=True)

    # Terminal panel (left)
    term_lines = [
        '$ claude "Create a Remotion video with',
        '  animated title and sound effects"',
        "",
        "> Creating project structure...",
        "> Writing Root.tsx...",
        "> Adding Composition with spring animation",
        "> Configuring audio track",
        "> Rendering video at 1920x1080, 30fps",
        "",
        "\u2714 Video rendered: output/video.mp4",
    ]
    add_code_block(slide, Inches(0.5), Inches(2.2), Inches(7.0), Inches(4.8),
                   term_lines, font_size=12, title="claude \u2014 terminal")

    # Files panel (right)
    files_l = Inches(8.0)
    files_t = Inches(2.2)
    files_w = Inches(4.8)
    files_h = Inches(4.8)
    add_rounded_rect(slide, files_l, files_t, files_w, files_h,
                     fill_color=DARK_PANEL,
                     border_color=RGBColor(0x33, 0x33, 0x33))

    add_textbox(slide, files_l + Inches(0.3), files_t + Inches(0.2),
                Inches(3), Inches(0.4), "PROJECT FILES",
                font_size=12, color=RGBColor(0x66, 0x66, 0x66), bold=True,
                alignment=PP_ALIGN.LEFT, rtl=False, font_name=FONT_MONO)

    files = [
        ("Root.tsx", "entry point"),
        ("Composition.tsx", "animated title"),
        ("AudioTrack.tsx", "sound effects"),
        ("package.json", "dependencies"),
    ]
    for i, (fname, fdesc) in enumerate(files):
        fy = files_t + Inches(0.8 + i * 0.9)
        # Green bullet
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     files_l + Inches(0.4), fy + Inches(0.08),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        # File name
        add_textbox(slide, files_l + Inches(0.7), fy - Inches(0.05),
                    Inches(3.5), Inches(0.35), fname,
                    font_size=14, color=WHITE, bold=True,
                    alignment=PP_ALIGN.LEFT, rtl=False, font_name=FONT_MONO)
        # Description
        add_textbox(slide, files_l + Inches(0.7), fy + Inches(0.25),
                    Inches(3.5), Inches(0.3), fdesc,
                    font_size=11, color=BODY_COLOR,
                    alignment=PP_ALIGN.LEFT, rtl=False, font_name=FONT_MONO)


def build_slide_8_skills(prs):
    """Slide 8: Skills install - progress bar"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_slide_label(slide, "\u05de\u05d4 \u05d6\u05d4 Skills")
    add_textbox(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.7),
                "\u05d7\u05d1\u05d9\u05dc\u05d5\u05ea \u05d9\u05db\u05d5\u05dc\u05d5\u05ea \u05de\u05d5\u05d1\u05e0\u05d5\u05ea \u05dc\u05e1\u05d5\u05db\u05df",
                font_size=30,
                color=HEADING_COLOR, bold=True,
                alignment=PP_ALIGN.CENTER, rtl=True)

    # Progress bar area
    bar_l = Inches(2.0)
    bar_w = Inches(9.3)
    bar_y = Inches(2.5)

    # Header: name + percentage
    add_textbox(slide, bar_l, bar_y, Inches(5), Inches(0.4),
                "Installing remotion-skill...",
                font_size=14, color=MUTED_COLOR, alignment=PP_ALIGN.LEFT,
                rtl=False, font_name=FONT_MONO)
    add_textbox(slide, bar_l + Inches(7), bar_y, Inches(2.3), Inches(0.4),
                "100%", font_size=14, color=ACCENT, bold=True,
                alignment=PP_ALIGN.RIGHT, rtl=False, font_name=FONT_MONO)

    # Track
    add_rounded_rect(slide, bar_l, bar_y + Inches(0.5), bar_w, Inches(0.3),
                     fill_color=RGBColor(0x33, 0x33, 0x33))
    # Fill (100%)
    add_rounded_rect(slide, bar_l, bar_y + Inches(0.5), bar_w, Inches(0.3),
                     fill_color=ACCENT)

    # Status text
    add_textbox(slide, bar_l, bar_y + Inches(0.95), bar_w, Inches(0.4),
                "\u2714 Loaded 847 tools \u00b7 \u2714 Ready to render videos",
                font_size=12, color=ACCENT, alignment=PP_ALIGN.LEFT,
                rtl=False, font_name=FONT_MONO)

    # Capability pills
    pills = ["React", "Animation", "Render", "Audio"]
    pill_w = Inches(1.8)
    pill_h = Inches(0.55)
    total_pills = len(pills) * pill_w + (len(pills) - 1) * Inches(0.3)
    pill_start = Emu(int((SLIDE_W - total_pills) / 2))

    for i, label in enumerate(pills):
        px = pill_start + i * (pill_w + Inches(0.3))
        add_pill(slide, px, Inches(4.2), pill_w, pill_h, label,
                 border_color=ACCENT, text_color=WHITE, font_size=14)

    # GitHub link text
    add_textbox(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(0.4),
                "remotion-dev/skills",
                font_size=13, color=BODY_COLOR, alignment=PP_ALIGN.CENTER,
                rtl=False, font_name=FONT_MONO)

    # Bottom note
    add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.5),
                "Skill = \u05de\u05d0\u05d5\u05ea \u05db\u05dc\u05d9\u05dd \u05e9\u05de\u05d0\u05e4\u05e9\u05e8\u05d9\u05dd \u05dc\u05e1\u05d5\u05db\u05df \u05dc\u05e2\u05d1\u05d5\u05d3 \u05e2\u05dd Remotion",
                font_size=15, color=BODY_COLOR, alignment=PP_ALIGN.CENTER,
                rtl=True)


def build_slide_9_checklist(prs):
    """Slide 9: Animated checklist (static, all checked)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_slide_label(slide, "\u05d0\u05d9\u05da \u05d6\u05d4 \u05d4\u05d5\u05dc\u05da \u05dc\u05e2\u05d1\u05d5\u05d3")
    add_textbox(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.7),
                "\u05de\u05d4 \u05ea\u05d3\u05e2\u05d5 \u05dc\u05e2\u05e9\u05d5\u05ea \u05e2\u05d3 \u05e1\u05d5\u05e3 \u05d4\u05e1\u05e8\u05d8\u05d5\u05df",
                font_size=30,
                color=HEADING_COLOR, bold=True,
                alignment=PP_ALIGN.CENTER, rtl=True)

    items = [
        "\u05d4\u05ea\u05e7\u05e0\u05ea Remotion",
        "\u05d9\u05e6\u05d9\u05e8\u05ea \u05e4\u05e8\u05d5\u05d9\u05e7\u05d8",
        "\u05d1\u05e0\u05d9\u05d9\u05ea \u05d0\u05e0\u05d9\u05de\u05e6\u05d9\u05d4",
        "\u05d4\u05d5\u05e1\u05e4\u05ea \u05e1\u05d0\u05d5\u05e0\u05d3",
        "\u05e8\u05d9\u05e0\u05d3\u05d5\u05e8 \u05d5\u05d9\u05e6\u05d5\u05d0 MP4",
    ]

    item_w = Inches(5)
    start_x = Emu(int((SLIDE_W - item_w) / 2))
    start_y = Inches(2.3)

    for i, text in enumerate(items):
        y = start_y + i * Inches(0.75)

        # Checkmark
        add_textbox(slide, start_x + Inches(3.8), y, Inches(0.5), Inches(0.5),
                    "\u2713", font_size=20, color=ACCENT, bold=True,
                    alignment=PP_ALIGN.CENTER, rtl=False)
        # Text
        add_textbox(slide, start_x, y, Inches(3.6), Inches(0.5),
                    text, font_size=20, color=MUTED_COLOR,
                    alignment=PP_ALIGN.RIGHT, rtl=True)

    # Bottom line
    add_textbox(slide, Inches(0), Inches(6.2), SLIDE_W, Inches(0.5),
                "\u05dc\u05d0 \u05ea\u05d9\u05d0\u05d5\u05e8\u05d9\u05d4. \u05d1\u05e0\u05d9\u05d9\u05d4 \u05d0\u05de\u05d9\u05ea\u05d9\u05ea.",
                font_size=20, color=ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER, rtl=True)


def build_slide_10_demo(prs):
    """Slide 10: Live demo transition"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0), Inches(1.0), SLIDE_W, Inches(0.5),
                "LIVE DEMO", font_size=16, color=ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER, rtl=False)

    # Two logos side by side with "x" between
    logo_y = Inches(2.0)
    logo_size = Inches(1.5)

    # Remotion logo placeholder
    add_rounded_rect(slide, Inches(4.0), logo_y, logo_size, logo_size,
                     fill_color=PANEL_BG, border_color=ACCENT)
    add_textbox(slide, Inches(4.0), logo_y + Inches(0.3), logo_size, Inches(0.8),
                "R", font_size=40, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, rtl=False)
    add_textbox(slide, Inches(3.5), logo_y + logo_size + Inches(0.1),
                Inches(2.5), Inches(0.4), "Remotion",
                font_size=16, color=MUTED_COLOR,
                alignment=PP_ALIGN.CENTER, rtl=False)

    # "x" between
    add_textbox(slide, Inches(5.9), logo_y + Inches(0.3), Inches(1.5), Inches(0.8),
                "\u00d7", font_size=36, color=BODY_COLOR,
                alignment=PP_ALIGN.CENTER, rtl=False)

    # Claude Code logo placeholder
    add_rounded_rect(slide, Inches(7.8), logo_y, logo_size, logo_size,
                     fill_color=PANEL_BG, border_color=ACCENT)
    add_textbox(slide, Inches(7.8), logo_y + Inches(0.3), logo_size, Inches(0.8),
                "A", font_size=40, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, rtl=False)
    add_textbox(slide, Inches(7.3), logo_y + logo_size + Inches(0.1),
                Inches(2.5), Inches(0.4), "Claude Code",
                font_size=16, color=MUTED_COLOR,
                alignment=PP_ALIGN.CENTER, rtl=False)

    # Main heading
    add_textbox(slide, Inches(0), Inches(4.5), SLIDE_W, Inches(0.8),
                "\u05e2\u05db\u05e9\u05d9\u05d5 \u05d1\u05d5\u05e0\u05d9\u05dd \u05d0\u05ea \u05d6\u05d4 \u05d1\u05d9\u05d7\u05d3",
                font_size=32, color=HEADING_COLOR, bold=True,
                alignment=PP_ALIGN.CENTER, rtl=True)

    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.5),
                "\u05de\u05e4\u05ea\u05d7\u05d9\u05dd \u05d0\u05ea \u05d4\u05e1\u05e8\u05d8\u05d5\u05df \u05d4\u05e8\u05d0\u05e9\u05d5\u05df \u05e9\u05dc\u05db\u05dd \u2014 \u05de\u05e4\u05e8\u05d5\u05de\u05e4\u05d8 \u05e2\u05d3 MP4",
                font_size=18, color=BODY_COLOR,
                alignment=PP_ALIGN.CENTER, rtl=True)

    # Down arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                   Inches(6.2), Inches(6.1),
                                   Inches(0.8), Inches(0.8))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT
    arrow.line.fill.background()


# ── Main ───────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1_title(prs)
    build_slide_2_code_to_video(prs)
    build_slide_3_problem(prs)
    build_slide_4_solution(prs)
    build_slide_5_showcase(prs)
    build_slide_6_pipeline(prs)
    build_slide_7_terminal(prs)
    build_slide_8_skills(prs)
    build_slide_9_checklist(prs)
    build_slide_10_demo(prs)

    output = "presentation.pptx"
    prs.save(output)
    print(f"Saved {output} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
